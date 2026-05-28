import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import plotly.express as px
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import io
import os

st.set_page_config(page_title="📊 Data Visualization Studio", layout="wide")

# 🎨 HEADER
st.markdown("""
    <div style="background: linear-gradient(90deg, #667eea 0%, #764ba2 100%); padding: 20px; border-radius: 10px; text-align: center;">
        <h1 style="color: white; margin: 0;">📊 Data Visualization Studio</h1>
        <p style="color: #ddd; margin: 0;">Create, visualize, and export your data insights</p>
    </div>
    """, unsafe_allow_html=True)

st.markdown("---")

# 📊 SESSION STATE
if 'df' not in st.session_state:
    st.session_state.df = None
if 'visualizations' not in st.session_state:
    st.session_state.visualizations = []

# 🔘 UPLOAD FILE
col1, col2, col3 = st.columns([2, 2, 2])

with col1:
    uploaded_file = st.file_uploader("📂 Load File", type=["csv", "xlsx", "xls"])
    if uploaded_file is not None:
        try:
            if uploaded_file.name.endswith('.csv'):
                st.session_state.df = pd.read_csv(uploaded_file)
            else:
                st.session_state.df = pd.read_excel(uploaded_file)
            st.success(f"✅ Loaded: {st.session_state.df.shape[0]} rows × {st.session_state.df.shape[1]} columns")
        except Exception as e:
            st.error(f"❌ Error: {str(e)}")

# 📑 TABS
if st.session_state.df is not None:
    tab1, tab2, tab3 = st.tabs(["📄 Data Preview", "📊 Create Charts", "⬇️ Export"])
    
    # TAB 1: DATA PREVIEW
    with tab1:
        st.subheader("Data Preview")
        st.dataframe(st.session_state.df, use_container_width=True)
        
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Rows", st.session_state.df.shape[0])
        with col2:
            st.metric("Columns", st.session_state.df.shape[1])
        with col3:
            st.metric("Memory", f"{st.session_state.df.memory_usage(deep=True).sum() / 1024:.2f} KB")
        
        st.subheader("Column Info")
        st.dataframe(pd.DataFrame({
            "Column": st.session_state.df.columns,
            "Type": st.session_state.df.dtypes.values,
            "Non-Null": st.session_state.df.count().values,
            "Null": st.session_state.df.isnull().sum().values
        }), use_container_width=True)
    
    # TAB 2: CREATE CHARTS
    with tab2:
        st.subheader("Create Charts")
        
        col1, col2 = st.columns(2)
        
        with col1:
            chart_type = st.selectbox(
                "📊 Chart Type",
                ["bar", "line", "pie", "histogram", "scatter", "box", "area"]
            )
        
        with col2:
            st.write("")  # spacing
        
        cols = st.session_state.df.columns.tolist()
        numeric_cols = st.session_state.df.select_dtypes(include=['number']).columns.tolist()
        
        col1, col2 = st.columns(2)
        
        with col1:
            x_col = st.selectbox("X Axis", cols)
        
        with col2:
            y_col = st.selectbox("Y Axis", numeric_cols if numeric_cols else cols)
        
        chart_title = st.text_input("Chart Title", value=f"{chart_type.capitalize()} Chart")
        
        if st.button("✨ Create Chart", type="primary"):
            try:
                fig = plt.figure(figsize=(10, 6))
                
                if chart_type == "bar":
                    plt.bar(st.session_state.df[x_col], st.session_state.df[y_col])
                elif chart_type == "line":
                    plt.plot(st.session_state.df[x_col], st.session_state.df[y_col], marker='o')
                elif chart_type == "pie":
                    plt.pie(st.session_state.df[y_col], labels=st.session_state.df[x_col], autopct='%1.1f%%')
                elif chart_type == "histogram":
                    plt.hist(st.session_state.df[x_col], bins=30, edgecolor='black')
                elif chart_type == "scatter":
                    plt.scatter(st.session_state.df[x_col], st.session_state.df[y_col], alpha=0.6)
                elif chart_type == "box":
                    plt.boxplot(st.session_state.df[y_col])
                elif chart_type == "area":
                    plt.fill_between(range(len(st.session_state.df)), st.session_state.df[y_col], alpha=0.3)
                    plt.plot(st.session_state.df[y_col])
                
                plt.title(chart_title, fontsize=16, fontweight='bold')
                plt.xlabel(x_col)
                plt.ylabel(y_col)
                plt.tight_layout()
                
                st.session_state.visualizations.append({
                    "title": chart_title,
                    "figure": fig,
                    "type": chart_type,
                    "x": x_col,
                    "y": y_col
                })
                
                st.success("✅ Chart created!")
                st.pyplot(fig)
                
            except Exception as e:
                st.error(f"❌ Error creating chart: {str(e)}")
        
        # Display saved visualizations
        if st.session_state.visualizations:
            st.markdown("---")
            st.subheader(f"📊 Saved Charts ({len(st.session_state.visualizations)})")
            
            for idx, viz in enumerate(st.session_state.visualizations):
                col1, col2 = st.columns([4, 1])
                with col1:
                    st.pyplot(viz["figure"])
                with col2:
                    if st.button("🗑️ Delete", key=f"del_{idx}"):
                        st.session_state.visualizations.pop(idx)
                        st.rerun()
    
    # TAB 3: EXPORT
    with tab3:
        if st.session_state.visualizations:
            st.subheader("⬇️ Export Options")
            
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("📄 Export as PDF", type="primary"):
                    try:
                        pdf_buffer = io.BytesIO()
                        with PdfPages(pdf_buffer) as pdf:
                            for viz in st.session_state.visualizations:
                                pdf.savefig(viz["figure"])
                        
                        pdf_buffer.seek(0)
                        st.download_button(
                            label="📥 Download PDF",
                            data=pdf_buffer.getvalue(),
                            file_name="visualizations.pdf",
                            mime="application/pdf"
                        )
                        st.success("✅ PDF ready for download!")
                    except Exception as e:
                        st.error(f"❌ Error: {str(e)}")
            
            with col2:
                if st.button("📽️ Export as PowerPoint", type="primary"):
                    try:
                        prs = Presentation()
                        
                        for viz in st.session_state.visualizations:
                            slide = prs.slides.add_slide(prs.slide_layouts[5])
                            
                            # Add title
                            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
                            title_frame = title_box.text_frame
                            title_frame.text = viz["title"]
                            title_frame.paragraphs[0].font.size = Inches(0.35)
                            title_frame.paragraphs[0].font.bold = True
                            
                            # Add image
                            img_buffer = io.BytesIO()
                            viz["figure"].savefig(img_buffer, format='png', bbox_inches='tight')
                            img_buffer.seek(0)
                            
                            slide.shapes.add_picture(img_buffer, Inches(0.5), Inches(1), width=Inches(9))
                        
                        ppt_buffer = io.BytesIO()
                        prs.save(ppt_buffer)
                        ppt_buffer.seek(0)
                        
                        st.download_button(
                            label="📥 Download PowerPoint",
                            data=ppt_buffer.getvalue(),
                            file_name="visualizations.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        st.success("✅ PowerPoint ready for download!")
                    except Exception as e:
                        st.error(f"❌ Error: {str(e)}")
        else:
            st.info("📊 Create some charts first to export!")

else:
    st.info("👆 Please upload a CSV or Excel file to get started!")
