import ttkbootstrap as ttk
from ttkbootstrap.constants import *
from tkinter import filedialog, messagebox
import tkinter as tk
import pandas as pd
import fitz  
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
from matplotlib.figure import Figure
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
import io
import os


class VisualizationPresenter:
    def __init__(self, root):
        self.root = root
        self.root.title("📊 Data Visualization Studio")
        self.root.geometry("1200x700")

        self.df = None
        self.pdf_text = ""
        self.visualizations = []

        self.setup_ui()

    def setup_ui(self):
        header = ttk.Frame(self.root, padding=15, bootstyle="dark")
        header.pack(fill=X)

        ttk.Label(
            header,
            text="📊 Data Visualization Studio",
            font=("Segoe UI", 20, "bold"),
            bootstyle="inverse-dark"
        ).pack(side=LEFT)

        top_frame = ttk.Frame(self.root, padding=10)
        top_frame.pack(fill=X)

        ttk.Button(
            top_frame,
            text="📂 Load File",
            bootstyle="primary",
            command=self.load_file
        ).pack(side=LEFT, padx=5)

        ttk.Button(
            top_frame,
            text="👀 Preview",
            bootstyle="info",
            command=self.preview_data
        ).pack(side=LEFT, padx=5)

        ttk.Button(
            top_frame,
            text="📊 Charts",
            bootstyle="success",
            command=self.create_visualizations_window
        ).pack(side=LEFT, padx=5)

        ttk.Button(
            top_frame,
            text="📽 PPT",
            bootstyle="warning",
            command=self.generate_ppt
        ).pack(side=LEFT, padx=5)

        ttk.Button(
            top_frame,
            text="📄 PDF",
            bootstyle="danger",
            command=self.export_pdf
        ).pack(side=LEFT, padx=5)

        self.file_label = ttk.Label(
            top_frame,
            text="No file loaded",
            bootstyle="secondary"
        )
        self.file_label.pack(side=LEFT, padx=20)

        self.notebook = ttk.Notebook(self.root, bootstyle="dark")
        self.notebook.pack(fill=BOTH, expand=True, padx=10, pady=10)

        self.preview_frame = ttk.Frame(self.notebook, padding=10)
        self.viz_frame = ttk.Frame(self.notebook, padding=10)

        self.notebook.add(self.preview_frame, text="📄 Data Preview")
        self.notebook.add(self.viz_frame, text="📊 Visualizations")

        self.setup_preview_tab()
        self.setup_viz_tab()

    def setup_preview_tab(self):
        self.preview_text = tk.Text(self.preview_frame, height=30)
        self.preview_text.pack(fill=BOTH, expand=True)

    def setup_viz_tab(self):
        self.viz_container = ttk.Frame(self.viz_frame)
        self.viz_container.pack(fill=BOTH, expand=True)

    def load_file(self):
        file_path = filedialog.askopenfilename(
            filetypes=[
                ("CSV Files", "*.csv"),
                ("Excel Files", "*.xlsx *.xls"),
                ("PDF Files", "*.pdf")
            ]
        )

        if not file_path:
            return

        try:
            ext = os.path.splitext(file_path)[1].lower()

            self.df = None
            self.pdf_text = ""

            if ext == ".csv":
                self.df = pd.read_csv(file_path)

            elif ext in [".xlsx", ".xls"]:
                self.df = pd.read_excel(file_path)

            elif ext == ".pdf":
                doc = fitz.open(file_path)

                text = ""
                for page in doc:
                    text += page.get_text()

                self.pdf_text = text

            else:
                raise ValueError("Unsupported file format")

            self.file_label.config(text=os.path.basename(file_path))

            if self.df is not None:
                messagebox.showinfo(
                    "Success",
                    f"Loaded Data File\nRows: {self.df.shape[0]}\nColumns: {self.df.shape[1]}"
                )
            else:
                messagebox.showinfo(
                    "Success",
                    "PDF Loaded Successfully"
                )

        except Exception as e:
            messagebox.showerror("Error", str(e))

    def preview_data(self):
        self.preview_text.delete("1.0", "end")

        if self.df is not None:
            self.preview_text.insert(
                "end",
                str(self.df.head(100))
            )

        elif self.pdf_text:
            self.preview_text.insert(
                "end",
                self.pdf_text[:50000]
            )

        else:
            messagebox.showwarning(
                "Warning",
                "No file loaded"
            )

    def create_visualizations_window(self):
        if self.df is None:
            messagebox.showwarning(
                "Warning",
                "Charts can only be created from CSV or Excel data."
            )
            return

        win = ttk.Toplevel(self.root)
        win.title("Create Chart")
        win.geometry("400x500")

        chart_var = ttk.StringVar(value="bar")

        for c in ["bar", "line", "pie", "histogram", "scatter", "box"]:
            ttk.Radiobutton(
                win,
                text=c,
                variable=chart_var,
                value=c
            ).pack(anchor="w")

        cols = self.df.columns.tolist()

        x_var = ttk.StringVar(value=cols[0])
        y_var = ttk.StringVar(value=cols[0])

        ttk.Label(win, text="X").pack()
        ttk.Combobox(win, textvariable=x_var, values=cols).pack()

        ttk.Label(win, text="Y").pack()
        ttk.Combobox(win, textvariable=y_var, values=cols).pack()

        title_entry = ttk.Entry(win)
        title_entry.pack(pady=10)

        def create_chart():
            fig = Figure(figsize=(6, 4))
            ax = fig.add_subplot(111)

            x = self.df[x_var.get()]
            y = self.df[y_var.get()]

            chart = chart_var.get()

            try:
                if chart == "bar":
                    ax.bar(x, y)

                elif chart == "line":
                    ax.plot(x, y)

                elif chart == "pie":
                    ax.pie(y, labels=x, autopct='%1.1f%%')

                elif chart == "histogram":
                    ax.hist(pd.to_numeric(x, errors="coerce").dropna())

                elif chart == "scatter":
                    ax.scatter(x, y)

                elif chart == "box":
                    ax.boxplot(pd.to_numeric(y, errors="coerce").dropna())

                title = title_entry.get() or "Chart"
                ax.set_title(title)

                self.visualizations.append((title, fig))
                self.display_visualizations()

                win.destroy()

            except Exception as e:
                messagebox.showerror("Chart Error", str(e))

        ttk.Button(
            win,
            text="Create Chart",
            bootstyle="success",
            command=create_chart
        ).pack(pady=10)

    def display_visualizations(self):
        for widget in self.viz_container.winfo_children():
            widget.destroy()

        for title, fig in self.visualizations:
            card = ttk.Frame(
                self.viz_container,
                padding=10,
                bootstyle="secondary"
            )
            card.pack(pady=10, fill=X)

            ttk.Label(
                card,
                text=title,
                font=("Segoe UI", 12, "bold")
            ).pack(anchor="w")

            canvas = FigureCanvasTkAgg(fig, master=card)
            canvas.draw()
            canvas.get_tk_widget().pack()

    def export_pdf(self):
        if not self.visualizations:
            messagebox.showwarning(
                "Warning",
                "No charts available."
            )
            return

        path = filedialog.asksaveasfilename(
            defaultextension=".pdf"
        )

        if not path:
            return

        with PdfPages(path) as pdf:
            for _, fig in self.visualizations:
                pdf.savefig(fig)

        messagebox.showinfo(
            "Success",
            "PDF Exported"
        )

    def generate_ppt(self):
        if not self.visualizations:
            messagebox.showwarning(
                "Warning",
                "No charts available."
            )
            return

        path = filedialog.asksaveasfilename(
            defaultextension=".pptx"
        )

        if not path:
            return

        prs = Presentation()

        for title, fig in self.visualizations:
            slide = prs.slides.add_slide(
                prs.slide_layouts[5]
            )

            textbox = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(0.3),
                Inches(5),
                Inches(0.5)
            )
            textbox.text = title

            img = io.BytesIO()
            fig.savefig(img, format='png')
            img.seek(0)

            slide.shapes.add_picture(
                img,
                Inches(1),
                Inches(1.2),
                width=Inches(6)
            )

        prs.save(path)

        messagebox.showinfo(
            "Success",
            "PPT Created"
        )


def main():
    root = ttk.Window(themename="darkly")
    app = VisualizationPresenter(root)
    root.mainloop()


if __name__ == "__main__":
    main()
